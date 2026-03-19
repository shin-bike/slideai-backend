"""
SlideAI Backend — FastAPI + python-pptx
Kimuraさんの手順通りに実装:
  1. 構成設計 (Claude API)
  2. テンプレート選択 (メタデータ検索)
  3. 空スライドを枚数分作成
  4. テンプレートスライドをコピーして適用
  5. テキスト流し込み
  6. 視覚QA (Claude Vision)
  7. 修正ループ
"""

import os, io, json, base64, re, copy, tempfile, shutil, logging
from pathlib import Path
from typing import Optional

from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, JSONResponse
from pydantic import BaseModel

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import lxml.etree as etree
import anthropic

logging.basicConfig(level=logging.INFO)
log = logging.getLogger("slideai")

BASE_DIR = Path(__file__).parent

# ── テンプレートファイルマップ ──
TEMPLATE_FILES = {
    "DTC_PowerLibrary_2013":  BASE_DIR / "DTC.pptx",
    "McKinsey":               BASE_DIR / "McKinsey.pptx",
    "BCG":                    BASE_DIR / "BCG.pptx",
    "PowerLibrary_2009":      BASE_DIR / "PowerLibrary2009.pptx",
}

# ── メタデータ読み込み ──
with open(BASE_DIR / "metadata.json", encoding="utf-8") as f:
    ALL_METADATA = json.load(f)

app = FastAPI(title="SlideAI API")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ════════════════════════════════════════
# リクエスト/レスポンス型
# ════════════════════════════════════════
class GenerateRequest(BaseModel):
    topic: str
    page_count: int = 5
    notes: str = ""
    api_key: str

class SlideSpec(BaseModel):
    page: int
    title: str
    purpose: str
    content_type: str
    key_points: list[str]
    visual_hint: str = ""
    headline: str = ""
    body: list[str] = []
    data_note: str = ""

# ════════════════════════════════════════
# Claude API ヘルパー
# ════════════════════════════════════════
def call_claude(api_key: str, prompt: str, system: str = "", max_tokens: int = 2000) -> str:
    client = anthropic.Anthropic(api_key=api_key)
    kwargs = dict(
        model="claude-sonnet-4-20250514",
        max_tokens=max_tokens,
        messages=[{"role": "user", "content": prompt}]
    )
    if system:
        kwargs["system"] = system
    msg = client.messages.create(**kwargs)
    return msg.content[0].text

def call_claude_vision(api_key: str, prompt: str, image_b64: str, media_type: str = "image/jpeg") -> str:
    client = anthropic.Anthropic(api_key=api_key)
    msg = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=1000,
        messages=[{
            "role": "user",
            "content": [
                {"type": "image", "source": {"type": "base64", "media_type": media_type, "data": image_b64}},
                {"type": "text", "text": prompt}
            ]
        }]
    )
    return msg.content[0].text

def parse_json(text: str):
    """JSONパース（途中切れや不正文字に強化）"""
    # コードブロック記法を除去
    cleaned = re.sub(r"```json\s*", "", text)
    cleaned = re.sub(r"```\s*", "", cleaned).strip()

    # まずそのままパース試行
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError as e:
        log.warning(f"JSON parse error: {e}. Attempting repair...")

    # 途中切れ対策: 不完全なJSONを修復
    # 最後の完全なオブジェクト（}）を見つけて配列を閉じる
    try:
        # 最後の "}," または "}" を見つけて配列として閉じる
        last_brace = cleaned.rfind('}')
        if last_brace > 0:
            truncated = cleaned[:last_brace+1]
            # 配列として閉じる
            if truncated.strip().startswith('['):
                truncated = truncated + ']'
            elif truncated.strip().startswith('{'):
                truncated = '[' + truncated + ']'
            result = json.loads(truncated)
            log.warning(f"JSON repaired: recovered {len(result)} items")
            return result
    except Exception as e2:
        log.warning(f"JSON repair also failed: {e2}")

    # それでも失敗したら元のエラーを再送出
    raise json.JSONDecodeError(f"Failed to parse JSON: {text[:200]}", text, 0)

# ════════════════════════════════════════
# Step 1: 構成設計
# ════════════════════════════════════════
def design_structure(api_key: str, topic: str, page_count: int, notes: str) -> list[dict]:
    prompt = f"""プレゼンテーション構成を設計してください。

テーマ: {topic}
ページ数: {page_count}
追加指示: {notes or 'なし'}

要件:
- 1ページ目は必ずタイトルスライド
- 最終ページはまとめ/結論スライド
- 各スライドのkey_pointsは5〜8個、各ポイントは30文字以上の具体的な内容
- 数値・事例・具体例を必ず含める
- content_typeは テキスト/データ/比較/プロセス/関係性/組織 から選ぶ

JSONのみ返答（コードブロック不要）:
[{{"page":1,"title":"スライドタイトル","purpose":"目的","content_type":"テキスト","key_points":["ポイント1（30文字以上）","ポイント2","ポイント3","ポイント4","ポイント5"],"visual_hint":"グラフ|テーブル|フロー|箇条書き|マトリクス"}}]"""

    tokens = max(2000, page_count * 300 + 500)
    raw = call_claude(api_key, prompt, max_tokens=tokens)
    return parse_json(raw)

# ════════════════════════════════════════
# Step 2: テンプレート選択
# ════════════════════════════════════════
# ════════════════════════════════════════
# Step 2: テンプレート選択
# ════════════════════════════════════════

# カテゴリ表紙・区切りスライドのタイトル（コンテンツ用途ではないため除外）
_EXCLUDE_TITLES = {
    'PowerPoint text tables', 'Text boxes', 'Chevrons', 'Charts', 'Graphs',
    'Shapes', 'Maps', 'Special Graphics', 'Qualifications', 'Resumes',
    'Numbered points', 'Major points', 'Data PowerPoint tables',
    'Alternative data PowerPoint tables', 'Gantt charts', 'Trees', 'Relations',
    'Structured Text', 'Power Library', 'Service Line Charts', 'Flags',
    'Consulting Report Template', 'Text Blocks', 'Org Charts', 'Driver Trees',
    'Logic Trees', 'Bars', 'Columns', 'Lines', 'Pies', 'Areas', 'Radars',
    'Waterfalls', 'Circles', 'Spheres', 'Triangles', 'Boxes', 'Puzzles',
    'Other Graphs', 'Combination line and bar chart', 'DTC 標準カラーパレット',
    'テンプレート検索ページ  1', 'テンプレート検索ページ  2',
    'Power Library & Appendix（日本語版）', 'Power Library（図形）',
    'Copy and paste these generic blocks, circles and arrows',
    'CONFIDENTIAL', 'END OF DOC', 'APPENDIX', 'STICKER AND OTHERS',
    'MAPS', 'Text', 'Label 1', 'Title', 'Update history', 'Table',
    'Chevrons with tables', 'Chevrons with text boxes', 'Chevrons with table and header column',
}

# 使用可能なメタデータ（カテゴリ表紙を除外済み）
_USABLE_METADATA = [m for m in ALL_METADATA if m.get('title', '').strip() not in _EXCLUDE_TITLES]

# DTC以外のメタデータ（top=None問題がないテンプレートを優先）
_NON_DTC_METADATA = [m for m in _USABLE_METADATA if m.get('source') != 'DTC_PowerLibrary_2013']

def select_template(purpose: str, content_type: str) -> Optional[dict]:
    """メタデータからベストマッチのテンプレートスライドを返す"""
    combined = (purpose + " " + content_type).lower()

    # DTC以外から優先して探す（top=None問題を回避）
    candidates = [m for m in _NON_DTC_METADATA if m.get("content_type") == content_type]
    if not candidates:
        candidates = _NON_DTC_METADATA
    if not candidates:
        candidates = _USABLE_METADATA

    # キーワードスコアリング
    def score(m):
        text = " ".join([
            m.get("title", ""),
            m.get("category", ""),
            m.get("use_case", ""),
            " ".join(m.get("keywords", []))
        ]).lower()
        return sum(1 for w in combined.split() if len(w) > 1 and w in text)

    ranked = sorted(candidates, key=score, reverse=True)
    return ranked[0] if ranked else None

# ════════════════════════════════════════
# Step 3: テンプレートスライドをコピー
# ════════════════════════════════════════
def copy_slide_from_template(src_prs: Presentation, slide_idx: int, dst_prs: Presentation) -> any:
    """テンプレートのスライドをdst_prsにコピーして返す"""
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    src_slide = src_prs.slides[slide_idx]

    # blank layoutでスライドを追加
    blank_layout = dst_prs.slide_layouts[6]
    new_slide = dst_prs.slides.add_slide(blank_layout)

    src_sp_tree = src_slide.shapes._spTree
    dst_sp_tree = new_slide.shapes._spTree

    # dst_sp_treeの既存要素を全削除
    for elem in list(dst_sp_tree):
        dst_sp_tree.remove(elem)

    # src_sp_treeの全要素をディープコピー
    for elem in src_sp_tree:
        dst_sp_tree.append(copy.deepcopy(elem))

    # 背景をコピー（XMLレベル）
    try:
        src_bg = src_slide._element.find(qn('p:bg'))
        if src_bg is not None:
            dst_slide_elem = new_slide._element
            existing_bg = dst_slide_elem.find(qn('p:bg'))
            if existing_bg is not None:
                dst_slide_elem.remove(existing_bg)
            sp_tree = dst_slide_elem.find('.//' + qn('p:spTree'))
            if sp_tree is not None:
                sp_tree.addprevious(copy.deepcopy(src_bg))
            else:
                dst_slide_elem.insert(0, copy.deepcopy(src_bg))
    except Exception as e:
        log.warning(f"Background copy failed: {e}")

    # プレースホルダーの位置情報を修正（top=Noneのシェイプに元スライドの座標を付与）
    src_shapes_by_name = {}
    for src_sh in src_slide.shapes:
        if src_sh.top is not None:
            src_shapes_by_name[src_sh.name] = src_sh

    for sh in new_slide.shapes:
        if sh.top is not None:
            continue
        if not sh.has_text_frame:
            continue
        src_sh = src_shapes_by_name.get(sh.name)
        if src_sh is None:
            continue
        try:
            sp_elem = sh._element
            spPr = sp_elem.find(qn('p:spPr'))
            if spPr is None:
                continue
            if spPr.find(qn('a:xfrm')) is not None:
                continue
            xfrm = etree.SubElement(spPr, qn('a:xfrm'))
            off  = etree.SubElement(xfrm, qn('a:off'))
            ext  = etree.SubElement(xfrm, qn('a:ext'))
            off.set('x', str(int(src_sh.left  or 0)))
            off.set('y', str(int(src_sh.top   or 0)))
            ext.set('cx', str(int(src_sh.width  or 9000000)))
            ext.set('cy', str(int(src_sh.height or 800000)))
        except Exception as e:
            log.warning(f"Position fix failed for {sh.name}: {e}")

    return new_slide

# ════════════════════════════════════════
# Step 4 & 5: テキスト流し込み
# ════════════════════════════════════════

# フッター・ページ番号の上端閾値（EMU）
_FOOTER_TOP = 6000000
# テンプレートの説明文パターン（除外対象）
_IGNORE_PATTERNS = [
    r'^横列項目', r'^シェブロンが不要', r'テキストボックスを縦位置',
    r'適宜削除', r'ドラッグ', r'^このスライド',
]

def _write_text(shape, lines: list) -> None:
    """XMLレベルでテキストを完全書き換え（PLACEHOLDER/TextBox両対応）"""
    try:
        txBody = shape.text_frame._txBody
        for p in txBody.findall(qn('a:p')):
            txBody.remove(p)
        for line in (lines if lines else [""]):
            p_elem = etree.SubElement(txBody, qn('a:p'))
            r_elem = etree.SubElement(p_elem, qn('a:r'))
            t_elem = etree.SubElement(r_elem, qn('a:t'))
            t_elem.text = str(line)
    except Exception as e:
        log.warning(f"_write_text error: {e}")

def _write_table(shape, body: list) -> None:
    """テーブルシェイプにbodyの内容をXMLレベルで流し込む"""
    try:
        tbl_elem = shape._element.find('.//' + qn('a:tbl'))
        if tbl_elem is None:
            return
        rows = tbl_elem.findall(qn('a:tr'))
        if not rows:
            return

        ncols = len(rows[0].findall(qn('a:tc')))

        # bodyを「見出し：内容」で分解
        def parse_item(item):
            sp = item.find('：')
            if sp > 0:
                return item[:sp].strip(), item[sp+1:].strip()
            return '', item.strip()

        body_idx = 0
        for ri, row in enumerate(rows):
            cells = row.findall(qn('a:tc'))
            for ci, cell in enumerate(cells):
                t_elems = cell.findall('.//' + qn('a:t'))
                existing = ''.join(t.text or '' for t in t_elems).strip()

                # プレースホルダーテキスト（全て書き換え対象）
                is_placeholder = existing in ('xx', '縦列タイトル', '横列項目', '') \
                    or 'テキスト' in existing or 'レベル' in existing

                if is_placeholder and body_idx < len(body):
                    key, val = parse_item(body[body_idx])
                    if ci == 0:
                        new_text = key or f'項目{body_idx+1}'
                    else:
                        new_text = val or body[body_idx]

                    # テキストを書き換え
                    for t_e in t_elems:
                        t_e.text = ''
                    if t_elems:
                        t_elems[0].text = new_text
                    else:
                        p_elem = cell.find('.//' + qn('a:p'))
                        if p_elem is None:
                            p_elem = etree.SubElement(cell, qn('a:p'))
                        r_elem = etree.SubElement(p_elem, qn('a:r'))
                        t_new  = etree.SubElement(r_elem, qn('a:t'))
                        t_new.text = new_text

                    # 最終列に達したら次のbody項目へ
                    if ci == ncols - 1:
                        body_idx += 1
    except Exception as e:
        log.warning(f"_write_table error: {e}")

def inject_content(slide, spec: dict) -> None:
    """スライドのテキストをXMLレベルで書き換える"""
    import re
    headline = spec.get("headline") or spec.get("title", "")
    body     = spec.get("body") or spec.get("key_points", [])

    title_shape    = None
    content_shapes = []
    table_shapes   = []
    seen_ids       = set()  # shape_idで重複を排除

    for s in slide.shapes:
        # shape_idで重複を排除
        sid = s.shape_id
        if sid in seen_ids:
            continue
        seen_ids.add(sid)

        top = s.top

        # フッター・ページ番号エリアを除外
        if top is not None and top > _FOOTER_TOP:
            continue

        # テーブルシェイプを収集（has_tableまたはXML内にa:tblがある場合）
        has_tbl = s.has_table or (s._element.find('.//' + qn('a:tbl')) is not None)
        if has_tbl:
            table_shapes.append(s)
            continue

        if not s.has_text_frame:
            continue

        text = s.text_frame.text.strip()

        # 数字のみはページ番号
        if text.isdigit():
            continue
        # テンプレート説明文を除外
        if any(re.search(p, text) for p in _IGNORE_PATTERNS):
            continue
        # 空のシェイプは除外（タイトル候補を除く）
        if not text and s.name != 'Rectangle 2':
            continue

        # Rectangle 2 = タイトル
        if s.name == 'Rectangle 2' and title_shape is None:
            title_shape = s
            continue

        content_shapes.append(s)

    # タイトルが見つからない場合は最初のcontentをタイトルに
    if not title_shape and content_shapes:
        title_shape = content_shapes.pop(0)

    # top順でソート
    content_shapes.sort(key=lambda s: (s.top if s.top is not None else 9999999, s.left or 0))

    # タイトルを書き換え
    if title_shape:
        _write_text(title_shape, [headline])

    # テーブルがある場合はそこにbodyを流し込む
    if table_shapes:
        _write_table(table_shapes[0], body)
        return

    if not content_shapes:
        return

    if len(content_shapes) == 1:
        _write_text(content_shapes[0], body)
    else:
        per = max(1, len(body) // len(content_shapes))
        for i, shape in enumerate(content_shapes):
            start = i * per
            end   = start + per if i < len(content_shapes) - 1 else len(body)
            chunk = body[start:end]
            _write_text(shape, chunk if chunk else [""])

def _set_text(shape, text: str, bold: bool = False) -> None:
    """後方互換用"""
    _write_text(shape, [text])

def _set_text_list(shape, items: list) -> None:
    """後方互換用"""
    _write_text(shape, items)

# ════════════════════════════════════════
# Step 5b: スライドを新規作成（デザインシステム使用）
# ════════════════════════════════════════
def create_slide_from_scratch(dst_prs: Presentation, spec: dict,
                               used_history: list = None) -> any:
    """slide_designs.py のデザイン関数を使ってスライドを生成"""
    from slide_designs import get_design_fn, DESIGN_NAMES

    headline = spec.get("headline") or spec.get("title", "")
    body     = spec.get("body") or spec.get("key_points", [])
    purpose  = spec.get("purpose", "")
    ct       = spec.get("content_type", "テキスト")
    page     = spec.get("page", 1)
    total    = spec.get("total_pages", 5)

    slide = dst_prs.slides.add_slide(dst_prs.slide_layouts[6])
    fn = get_design_fn(purpose, ct, page, total, used_history or [])
    fn(slide, dst_prs, headline, body)

    # 使用履歴を更新
    if used_history is not None:
        used_history.append(DESIGN_NAMES.get(fn, 'detail'))

    return slide

def _legacy_create_slide(dst_prs, spec):
    """旧実装（未使用）"""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    slide_layout = dst_prs.slide_layouts[6]  # Blank
    slide = dst_prs.slides.add_slide(slide_layout)

    W = int(dst_prs.slide_width)
    H = int(dst_prs.slide_height)
    MX = int(Inches(0.45))
    MY = int(Inches(0.35))
    TH = int(Inches(0.75))
    CY = int(MY + TH + int(Inches(0.15)))
    CW = int(W - MX * 2)
    CH = int(H - CY - int(Inches(0.25)))

    headline = spec.get("headline") or spec.get("title", "")
    body     = spec.get("body") or spec.get("key_points", [])
    ct       = spec.get("content_type", "テキスト")

    NAVY  = RGBColor(0x1E, 0x3A, 0x5F)
    BLUE  = RGBColor(0x25, 0x63, 0xB0)
    LGRAY = RGBColor(0xCA, 0xCE, 0xD8)
    LIGHT = RGBColor(0xEE, 0xF3, 0xFB)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY  = RGBColor(0x5A, 0x64, 0x78)

    def add_rect(x, y, w, h, color):
        shape = slide.shapes.add_shape(1, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_textbox(x, y, w, h, text, size=12, bold=False, color=None, wrap=True, align=PP_ALIGN.LEFT):
        txb = slide.shapes.add_textbox(x, y, w, h)
        tf  = txb.text_frame
        tf.word_wrap = wrap
        tf.auto_size = None
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or NAVY
        return txb

    # タイトル（全スライド共通）
    add_textbox(MX, MY, CW, TH, headline, size=22, bold=True, color=NAVY)
    add_rect(MX, int(MY + TH), CW, 25000, LGRAY)

    # ヘルパー: Inches を int に変換
    def I(n): return int(Inches(n))

    if ct == "プロセス":
        steps = body[:5]
        n = len(steps) or 3
        step_w = int((CW - I(0.18) * (n-1)) / n)
        accent_colors = [BLUE, RGBColor(0x1A,0x7A,0x52), RGBColor(0xD4,0x88,0x0A), RGBColor(0x4B,0x8E,0xE8), RGBColor(0xC0,0x39,0x2B)]
        for i, step in enumerate(steps):
            sx = int(MX + i * (step_w + I(0.18)))
            ac = accent_colors[i % len(accent_colors)]
            add_rect(sx, CY, step_w, CH, LIGHT)
            add_rect(sx, CY, step_w, 70000, ac)
            add_textbox(sx + I(0.15), CY + I(0.15), step_w - I(0.3), I(0.5),
                        f"0{i+1}", size=26, bold=True, color=ac)
            sp = step.find("：")
            head = step[:sp] if sp > 0 else step[:min(len(step), 18)]
            detail = step[sp+1:] if sp > 0 else step[18:]
            add_textbox(sx + I(0.15), CY + I(0.8), step_w - I(0.3), I(0.65),
                        head, size=12, bold=True, color=NAVY)
            if detail:
                add_textbox(sx + I(0.15), CY + I(1.55), step_w - I(0.3), max(I(0.3), CH - I(1.8)),
                            detail, size=10, color=GRAY)

    elif ct == "比較":
        cols = min(len(body), 4)
        label_w = I(1.5)
        cell_w  = int((CW - label_w) / max(cols, 1))
        rows = ["自社", "競合A", "競合B"]
        row_h = int(CH / (len(rows) + 1))
        # ヘッダー
        add_rect(MX, CY, CW, row_h, NAVY)
        add_textbox(MX, CY, label_w, row_h, "評価軸", size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for ci, label in enumerate(body[:cols]):
            add_textbox(MX + label_w + ci * cell_w, CY, cell_w, row_h,
                        label, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # 行
        eval_data = [
            ["◎ 特化・高品質", "◎ 競争力高い", "◎ 独自差別化", "◎ ブランド力"],
            ["△ 差別化薄い", "○ 標準品質", "△ 革新性低い", "○ 認知度高い"],
            ["○ 低価格訴求", "△ 品質ばらつき", "○ 立地優位", "△ サービス浅い"],
        ]
        for ri, row_label in enumerate(rows):
            ry = int(CY + (ri + 1) * row_h)
            fill = LIGHT if ri % 2 == 0 else WHITE
            add_rect(MX, ry, CW, row_h, fill)
            add_rect(MX, ry, label_w, row_h, BLUE if ri == 0 else NAVY)
            add_textbox(MX, ry, label_w, row_h, row_label, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            for ci in range(cols):
                eval_text = eval_data[ri][ci] if ci < len(eval_data[ri]) else "－"
                add_textbox(int(MX + label_w + ci * cell_w + I(0.08)), int(ry + I(0.06)),
                            int(cell_w - I(0.16)), int(row_h - I(0.12)),
                            eval_text, size=10, color=NAVY if ri == 0 else GRAY, align=PP_ALIGN.CENTER)

    else:
        # デフォルト: 2列テキスト
        pts  = body[:8]
        half = (len(pts) + 1) // 2
        col_w = int((CW - I(0.2)) / 2)
        item_h = int(min(CH / max(half, 1) - I(0.08), I(1.4)))
        for col_idx, group in enumerate([pts[:half], pts[half:]]):
            gx = int(MX + col_idx * (col_w + I(0.2)))
            for i, pt in enumerate(group):
                iy = int(CY + i * (item_h + I(0.08)))
                ac = [BLUE, RGBColor(0x1A,0x7A,0x52), RGBColor(0xD4,0x88,0x0A), RGBColor(0x4B,0x8E,0xE8)][i % 4]
                add_rect(gx, iy, col_w, item_h, LIGHT)
                add_rect(gx, iy, I(0.05), item_h, ac)
                sp = pt.find("：")
                if sp > 0:
                    add_textbox(gx + I(0.12), iy + I(0.06), col_w - I(0.2), I(0.32),
                                pt[:sp], size=11, bold=True, color=NAVY)
                    add_textbox(gx + I(0.12), iy + I(0.40), col_w - I(0.2), max(I(0.2), item_h - I(0.5)),
                                pt[sp+1:], size=10, color=GRAY)
                else:
                    add_textbox(gx + I(0.12), iy + I(0.08), col_w - I(0.2), max(I(0.2), item_h - I(0.16)),
                                pt, size=11, color=RGBColor(0x1A, 0x2A, 0x44))

    return slide

# ════════════════════════════════════════
# Step 6: 視覚QA (Claude Vision)
# ════════════════════════════════════════
def qa_slides_with_vision(api_key: str, prs: Presentation, tmp_dir: Path) -> list[dict]:
    """各スライドをPNG化してClaude Visionでチェック"""
    issues = []

    # LibreOfficeでPDF変換
    pptx_path = tmp_dir / "check.pptx"
    pdf_path  = tmp_dir / "check.pdf"
    prs.save(str(pptx_path))

    ret = os.system(f"libreoffice --headless --convert-to pdf '{pptx_path}' --outdir '{tmp_dir}' 2>/dev/null")
    if ret != 0 or not pdf_path.exists():
        log.warning("PDF conversion failed, skipping QA")
        return []

    # PDFを画像に変換
    ret = os.system(f"pdftoppm -jpeg -r 100 '{pdf_path}' '{tmp_dir}/slide' 2>/dev/null")
    if ret != 0:
        log.warning("pdftoppm failed, skipping QA")
        return []

    slide_images = sorted(tmp_dir.glob("slide-*.jpg"))
    if not slide_images:
        slide_images = sorted(tmp_dir.glob("slide*.jpg"))

    for i, img_path in enumerate(slide_images):
        with open(img_path, "rb") as f:
            img_b64 = base64.b64encode(f.read()).decode()

        prompt = """このスライドを視覚的に確認してください。
以下の問題があれば報告してください（なければ「問題なし」と返答）:
- テキストがボックスからはみ出している
- テキストが非常に小さくて読めない
- 要素が重なっている
- 空のプレースホルダーが残っている
- 明らかなレイアウト崩れ

問題があれば「問題あり: [内容]」、なければ「問題なし」のみ返答。"""

        try:
            result = call_claude_vision(api_key, prompt, img_b64)
            if "問題あり" in result:
                issues.append({"slide_index": i, "issue": result})
                log.info(f"Slide {i+1} QA issue: {result}")
        except Exception as e:
            log.warning(f"Vision QA failed for slide {i+1}: {e}")

    return issues

# ════════════════════════════════════════
# Step 7: コンテンツ生成
# ════════════════════════════════════════
def generate_content(api_key: str, plan: list[dict], topic: str, notes: str) -> list[dict]:
    """ページ数に応じて分割生成してJSONエラーを防ぐ"""
    CHUNK_SIZE = 5  # 一度に生成するページ数の上限

    all_results = []
    chunks = [plan[i:i+CHUNK_SIZE] for i in range(0, len(plan), CHUNK_SIZE)]

    for chunk in chunks:
        prompt = f"""各スライドの詳細コンテンツを生成してください。

テーマ: {topic}
追加指示: {notes or 'なし'}

対象スライド（{len(chunk)}ページ分）:
{chr(10).join(f"Page{s['page']}: {s['title']}（{s['purpose']}）" for s in chunk)}

条件:
- headlineは20〜45文字のインパクトある1文（数値・体言止めを使う）
- bodyは6〜8個のポイント、「見出し：詳細説明」形式（各40文字以上）
- 具体的な数値・事例・根拠を必ず含める
- data_noteはデータスライドのみ「65,78,85,79,92,98」形式で数値を記載
- 必ず全{len(chunk)}ページ分を生成すること

JSONのみ返答（コードブロック不要）:
[{{"page":1,"headline":"メインメッセージ","body":["見出し1：詳細説明40文字以上","見出し2：...","見出し3：...","見出し4：...","見出し5：...","見出し6：..."],"data_note":""}}]"""

        # チャンクサイズに応じてmax_tokensを動的に設定（1ページ=700トークン見込み）
        tokens = max(2000, len(chunk) * 700 + 500)
        raw = call_claude(api_key, prompt, max_tokens=tokens)
        chunk_results = parse_json(raw)
        all_results.extend(chunk_results)

    return all_results

# ════════════════════════════════════════
# メイン生成エンドポイント
# ════════════════════════════════════════
@app.post("/generate")
async def generate(req: GenerateRequest):
    log.info(f"Generate request: topic='{req.topic}', pages={req.page_count}")
    tmp_dir = Path(tempfile.mkdtemp())

    try:
        # Step 1: 構成設計
        log.info("Step 1: Designing structure...")
        plan = design_structure(req.api_key, req.topic, req.page_count, req.notes)
        log.info(f"  → {len(plan)} slides planned")

        # Step 3: コンテンツ生成
        log.info("Step 3: Generating content...")
        contents = generate_content(req.api_key, plan, req.topic, req.notes)
        content_map = {c["page"]: c for c in contents}

        # 出力pptxを作成（スライドサイズはDTCテンプレートに合わせる）
        ref_prs = Presentation(str(TEMPLATE_FILES["DTC_PowerLibrary_2013"]))
        dst_prs = Presentation()
        dst_prs.slide_width  = ref_prs.slide_width   # 10.83"
        dst_prs.slide_height = ref_prs.slide_height  # 7.5"

        # テンプレートファイルをキャッシュ
        loaded_templates = {}

        slide_info = []
        used_history = []  # デザイン使用履歴（重複回避用）

        for slide_spec in plan:
            page = slide_spec["page"]
            ct   = slide_spec.get("content_type", "テキスト")
            spec = {**slide_spec, **(content_map.get(page, {})), "total_pages": len(plan)}

            log.info(f"Creating slide {page} ({ct}) from scratch...")
            create_slide_from_scratch(dst_prs, spec, used_history)
            slide_info.append({"page": page, "used_template": False})

        # Step 5: 一時保存
        pptx_path = tmp_dir / "output.pptx"
        dst_prs.save(str(pptx_path))

        # Step 6: 視覚QA
        log.info("Step 6: Visual QA...")
        qa_issues = qa_slides_with_vision(req.api_key, dst_prs, tmp_dir)
        if qa_issues:
            log.info(f"  → {len(qa_issues)} issues found (logged, auto-fix not yet implemented)")

        # Step 7: 完了 → バイナリ返却
        with open(pptx_path, "rb") as f:
            pptx_bytes = f.read()

        import datetime, urllib.parse
        date_str = datetime.date.today().isoformat()
        # ASCII safe filename for compatibility
        fname_ascii = f"presentation_{date_str}.pptx"
        # RFC 5987 encoded filename for Japanese support
        fname_utf8 = urllib.parse.quote(f"{req.topic[:20]}_{date_str}.pptx")
        return StreamingResponse(
            io.BytesIO(pptx_bytes),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f"attachment; filename=\"{fname_ascii}\"; filename*=UTF-8''{fname_utf8}",
                "X-Slide-Count": str(len(plan)),
                "X-QA-Issues": str(len(qa_issues)),
            }
        )

    except Exception as e:
        log.error(f"Generation failed: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


@app.get("/health")
def health():
    return {"status": "ok", "templates": {k: v.exists() for k, v in TEMPLATE_FILES.items()}}

@app.get("/")
def root():
    return {"message": "SlideAI API", "docs": "/docs"}
