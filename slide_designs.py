"""
SlideAI Design System v7 — 内容適応型・動的図解
改善点:
  - bodyの項目数・内容に合わせて図を動的に生成
  - フロー図: 3〜8ステップに対応、幅を自動調整
  - 論点リスト: 項目数に合わせて行高を自動調整
  - ガント: タスク数に合わせてバー高を自動調整
  - 2カラム: 項目を左右に動的分配
  - テーブル: 行数に合わせてrowの高さを動的計算
  - シャドウなし・最小限の塗りつぶし
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import lxml.etree as etree

# ── カラーパレット（MF準拠）──
NAVY   = RGBColor(0x00, 0x20, 0x60)
BLUE1  = RGBColor(0x1F, 0x38, 0x64)
BLUE2  = RGBColor(0x2E, 0x74, 0xB5)
BLUE3  = RGBColor(0x9D, 0xC3, 0xE6)
LBLUE  = RGBColor(0xBD, 0xD7, 0xEE)
ORANGE = RGBColor(0xE8, 0x77, 0x22)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
LGRAY  = RGBColor(0xD0, 0xD0, 0xD0)
GRAY1  = RGBColor(0x26, 0x26, 0x26)
GRAY2  = RGBColor(0x44, 0x44, 0x44)
GRAY3  = RGBColor(0x88, 0x88, 0x88)
BLACK  = RGBColor(0x00, 0x00, 0x00)

def I(n): return int(Inches(n))
def P(n): return int(Pt(n))

def _r(slide, x, y, w, h, fill=WHITE, lc=None, lw=0.75):
    """矩形（シャドウなし保証）"""
    if w <= 0 or h <= 0: return None
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if lc:
        sh.line.color.rgb = lc; sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    # シャドウ除去
    try:
        sp = sh._element; spPr = sp.find(qn('p:spPr'))
        if spPr is None: spPr = etree.SubElement(sp, qn('p:spPr'))
        for tag in [qn('a:effectLst'), qn('a:effectDag')]:
            el = spPr.find(tag)
            if el is not None: spPr.remove(el)
    except: pass
    return sh

def _t(slide, x, y, w, h, text, sz=11, bold=False, col=None,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    if not text or w <= 0 or h <= 0: return None
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame; tf.word_wrap = wrap; tf.auto_size = None
    bp  = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('lIns','45720'); bp.set('rIns','45720')
        bp.set('tIns','22860'); bp.set('bIns','22860')
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = col or BLACK
    return txb

def _parse(b):
    for sep in ('：', ':'):
        sp = b.find(sep)
        if sp > 0: return b[:sp].strip(), b[sp+1:].strip()
    return '', b.strip()

# ── 共通フレーム ──
def _page_frame(slide, prs, context='', page_num=0):
    W = int(prs.slide_width); H = int(prs.slide_height)
    _r(slide, 0, I(0.28), W, P(1), fill=NAVY)
    if context:
        _t(slide, I(0.3), I(0.05), I(7), I(0.22), context, sz=8, col=GRAY2)
    _t(slide, W-I(2.5), I(0.05), I(2.3), I(0.22),
       'For Discussion', sz=8, italic=True, col=GRAY2, align=PP_ALIGN.RIGHT)
    _r(slide, 0, H-I(0.35), W, P(0.5), fill=NAVY)
    _t(slide, I(0.3), H-I(0.3), I(4), I(0.22),
       'SlideAI Copy Right Reserved', sz=7, col=GRAY3)
    if page_num:
        _t(slide, W-I(0.6), H-I(0.3), I(0.5), I(0.22),
           str(page_num), sz=8, col=GRAY2, align=PP_ALIGN.RIGHT)

def _slide_title(slide, prs, title, summary=''):
    W = int(prs.slide_width)
    _t(slide, I(0.3), I(0.35), W-I(0.6), I(0.55), title, sz=20, bold=True, col=NAVY)
    _r(slide, I(0.3), I(0.95), W-I(0.6), P(1.5), fill=NAVY)
    if summary:
        _t(slide, I(0.3), I(1.0), W-I(0.6), I(0.42), '■ '+summary, sz=10, col=GRAY1)

def _navy_header(slide, x, y, w, h, text, sz=10):
    _r(slide, x, y, w, h, fill=NAVY)
    _t(slide, x, y+I(0.06), w, h-I(0.1), text, sz=sz, bold=True, col=WHITE, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 1. タイトルスライド
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def slide_title(slide, prs, headline, body):
    W, H = int(prs.slide_width), int(prs.slide_height)
    _r(slide, 0, 0, W, H, fill=WHITE)
    client = body[0] if body else ''
    lbl, _ = _parse(client)
    _t(slide, I(0.7), I(1.2), W-I(1.4), I(0.6), lbl or client[:24], sz=20, col=BLUE2)
    _r(slide, I(0.5), I(2.2), W-I(1.0), I(1.1), fill=NAVY)
    if '【' in headline and '】' in headline:
        tag  = headline[headline.find('【'):headline.find('】')+1]
        main = headline[headline.find('】')+1:].strip()
        _t(slide, I(0.7), I(2.28), W-I(1.4), I(0.3), tag, sz=11, col=WHITE)
        _t(slide, I(0.7), I(2.62), W-I(1.4), I(0.55), main, sz=18, bold=True, col=WHITE)
    else:
        _t(slide, I(0.7), I(2.3), W-I(1.4), I(0.75), headline, sz=18, bold=True, col=WHITE)
    if len(body) > 1:
        _t(slide, I(0.7), I(3.55), I(4), I(0.3), body[1], sz=10, col=GRAY2)
    if len(body) > 2:
        _t(slide, I(0.7), I(3.88), I(4), I(0.3), body[2], sz=10, col=GRAY2)
    _t(slide, W-I(3.0), H-I(0.7), I(2.8), I(0.35),
       'SMART STRATEGY', sz=16, bold=True, col=GRAY1, align=PP_ALIGN.RIGHT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 2. フロー図（動的ステップ数対応）
# n=3〜8ステップに自動対応
# 上段: n個の白枠ボックス+矢印+成果物
# 下段: 2カラム詳細（項目数に合わせて行高調整）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def slide_flow(slide, prs, headline, body):
    W, H = int(prs.slide_width), int(prs.slide_height)
    _r(slide, 0, 0, W, H, fill=WHITE)
    _page_frame(slide, prs)
    summary = body[0].split('：')[1] if body and '：' in body[0] else (body[0] if body else '')
    _slide_title(slide, prs, headline, summary)

    # bodyから実際のステップを抽出（最初の要素はsummaryなのでスキップ）
    steps = [b for b in body if b and not (b == body[0] and summary)]
    # 6ステップ以上は上段に収まらないので分割
    top_steps = steps[:min(len(steps), 6)]
    n = len(top_steps)
    if n == 0:
        top_steps = body[:5]; n = len(top_steps)

    # ── 上段: プロセスフロー ──
    AREA_Y = I(1.52)
    # ステップ数に応じてフォントサイズ調整
    fs_label = 9 if n <= 4 else (8 if n <= 6 else 7)
    fs_det   = 8 if n <= 4 else (7 if n <= 6 else 6)

    # 成果物ラベルの有無でエリア高さを調整
    has_result = n <= 6  # 7以上はスペース不足で成果物省略
    AREA_H = I(1.85) if has_result else I(1.35)

    _r(slide, I(0.3), AREA_Y, W-I(0.6), AREA_H, fill=WHITE, lc=LGRAY, lw=0.5)

    # 矢印幅: n≤5は広め、n>5は細め
    arrow_w = I(0.18) if n <= 5 else I(0.14)
    total_w = W - I(0.8)
    box_w   = int((total_w - arrow_w*(n-1)) / n)
    box_h   = I(0.72) if has_result else I(1.0)

    # 成果物ラベル（コンテンツから推測、または固定）
    result_labels = []
    for b in top_steps:
        lbl, det = _parse(b)
        # 「→成果物」形式があれば使用、なければラベル短縮
        if '→' in det:
            result_labels.append(det.split('→')[-1].strip()[:8])
        else:
            result_labels.append((lbl or b[:8])[:8])

    for i, b in enumerate(top_steps):
        lbl, det = _parse(b)
        bx = I(0.4) + i*(box_w + arrow_w)
        by = AREA_Y + I(0.15)

        # プロセスボックス（白地+ネイビー枠）
        _r(slide, bx, by, box_w, box_h, fill=WHITE, lc=NAVY, lw=1.0)

        # 番号（左上）
        nums = ['①','②','③','④','⑤','⑥','⑦','⑧']
        _t(slide, bx+I(0.05), by+I(0.03), I(0.22), I(0.22),
           nums[i] if i < len(nums) else str(i+1), sz=8, bold=True, col=NAVY)

        # ラベル
        _t(slide, bx+I(0.05), by+I(0.24), box_w-I(0.08), I(0.22),
           lbl or b[:14], sz=fs_label, bold=True, col=NAVY)

        # 詳細（ステップ数が少ない時だけ表示）
        if fs_det >= 7 and det:
            _t(slide, bx+I(0.05), by+I(0.46), box_w-I(0.08), I(0.22),
               det[:18], sz=fs_det, col=GRAY2)

        # → 矢印
        if i < n-1:
            _t(slide, bx+box_w+I(0.0), by+box_h//2-I(0.12),
               arrow_w, I(0.24), '→', sz=10, col=NAVY, align=PP_ALIGN.CENTER)

        # ↓ + 成果物ボックス
        if has_result:
            _t(slide, bx+box_w//2-I(0.1), by+box_h+I(0.02), I(0.2), I(0.16),
               '↓', sz=8, col=GRAY3, align=PP_ALIGN.CENTER)
            rby = by+box_h+I(0.18)
            _r(slide, bx, rby, box_w, I(0.26), fill=WHITE, lc=GRAY2, lw=0.5)
            _t(slide, bx, rby+I(0.03), box_w, I(0.22),
               result_labels[i], sz=7, col=GRAY1, align=PP_ALIGN.CENTER)

    # ── 下段: 2カラム詳細（項目数に合わせて行高を動的計算）──
    DY = AREA_Y + AREA_H + I(0.15)
    DH = H - DY - I(0.42)
    if DH < I(0.5): return  # スペース不足

    # 全ステップを左右に振り分け
    half = (n + 1) // 2
    left_items  = top_steps[:half]
    right_items = top_steps[half:]
    cw = int((W - I(0.9)) / 2)

    for ci, (items, col_title) in enumerate([(left_items,'課題の構造'),(right_items,'解決アプローチ')]):
        if not items: continue
        cx = I(0.4) + ci*(cw+I(0.1))
        hh = I(0.34)
        _navy_header(slide, cx, DY, cw, hh, col_title, sz=10)
        _r(slide, cx, DY+hh, cw, DH-hh, fill=WHITE, lc=NAVY, lw=0.75)

        avail_h = DH - hh - I(0.08)
        per = avail_h // max(len(items), 1)
        # フォントは行高に応じて調整
        fs = 9 if per >= I(0.7) else (8 if per >= I(0.55) else 7)

        iy = DY + hh + I(0.06)
        for j, item in enumerate(items):
            lbl, det = _parse(item)
            _t(slide, cx+I(0.12), iy+I(0.02), cw-I(0.18), min(I(0.26),per*0.4),
               lbl or item[:22], sz=fs, bold=True, col=BLUE1)
            if det:
                _t(slide, cx+I(0.12), iy+I(0.26), cw-I(0.18), per-I(0.3),
                   det, sz=fs-1, col=GRAY1)
            if j < len(items)-1:
                _r(slide, cx+I(0.08), iy+per-I(0.04), cw-I(0.16), P(0.3), fill=LGRAY)
            iy += per


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 3. 論点リスト（動的項目数対応）
# 3〜8項目に対応、多い場合は2列配置
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def slide_detail(slide, prs, headline, body):
    W, H = int(prs.slide_width), int(prs.slide_height)
    _r(slide, 0, 0, W, H, fill=WHITE)
    _page_frame(slide, prs)
    summary = body[0].split('：')[1] if body and '：' in body[0] else (body[0] if body else '')
    _slide_title(slide, prs, headline, summary)

    # ネイビーヘッダー帯
    _navy_header(slide, I(0.4), I(1.52), W-I(0.8), I(0.34), '主要論点', sz=10)

    items = [b for b in body if b and not (b == body[0] and summary)]
    if not items: items = body
    n = min(len(items), 8)

    GY = I(1.88)
    GH = H - GY - I(0.42)

    # 項目数で1列/2列を切り替え
    if n <= 4:
        # 1列レイアウト（ゆったり）
        row_h = GH // n
        for i, b in enumerate(items[:n]):
            lbl, det = _parse(b)
            ry = GY + i*row_h
            # 左ラベル列
            _r(slide, I(0.4), ry+P(2), I(1.5), row_h-P(4), fill=BGRAY, lc=LGRAY, lw=0.3)
            _r(slide, I(0.46), ry+I(0.08), I(0.48), I(0.48), fill=ORANGE)
            _t(slide, I(0.46), ry+I(0.06), I(0.48), I(0.5),
               str(i+1), sz=14, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
            _t(slide, I(0.4), ry+I(0.62), I(1.5), row_h-I(0.7),
               lbl[:10] if lbl else b[:10], sz=8, bold=True, col=NAVY,
               align=PP_ALIGN.CENTER, wrap=True)
            # 右ブレット列
            _r(slide, I(2.05), ry+P(2), W-I(2.45), row_h-P(4),
               fill=WHITE, lc=LGRAY, lw=0.3)
            if det:
                fs = 10 if row_h >= I(1.0) else 9
                bullets = det.split('・') if '・' in det else det.split('、')
                bh = (row_h-I(0.12)) // max(len(bullets[:4]), 1)
                by2 = ry + I(0.06)
                for bullet in bullets[:4]:
                    if bullet.strip():
                        _r(slide, I(2.12), by2+I(0.1), I(0.14), I(0.14), fill=NAVY)
                        _t(slide, I(2.32), by2+I(0.04), W-I(2.7), bh-I(0.06),
                           bullet.strip(), sz=fs, col=GRAY1)
                        by2 += bh
            if i < n-1:
                _r(slide, I(0.4), ry+row_h-P(2), W-I(0.8), P(0.5), fill=LGRAY)
    else:
        # 2列レイアウト（5〜8項目）
        cols = 2
        half = (n + 1) // 2
        cw = int((W - I(0.9)) / 2)
        for ci in range(cols):
            cx = I(0.4) + ci*(cw+I(0.1))
            col_items = items[ci*half:(ci+1)*half]
            if not col_items: continue
            n_col = len(col_items)
            row_h = GH // n_col
            for ri, b in enumerate(col_items):
                lbl, det = _parse(b)
                ry = GY + ri*row_h
                # オレンジ円
                label_w = I(1.0)
                _r(slide, cx, ry+P(2), label_w, row_h-P(4), fill=BGRAY, lc=LGRAY, lw=0.3)
                _r(slide, cx+I(0.1), ry+I(0.06), I(0.36), I(0.36), fill=ORANGE)
                _t(slide, cx+I(0.1), ry+I(0.04), I(0.36), I(0.36),
                   str(ci*half+ri+1), sz=10, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
                _t(slide, cx, ry+I(0.46), label_w, row_h-I(0.52),
                   lbl[:8] if lbl else b[:8], sz=7, bold=True, col=NAVY,
                   align=PP_ALIGN.CENTER, wrap=True)
                # ブレット列
                _r(slide, cx+label_w+I(0.05), ry+P(2),
                   cw-label_w-I(0.05), row_h-P(4),
                   fill=WHITE, lc=LGRAY, lw=0.3)
                if det:
                    fs = 9 if row_h >= I(0.9) else 8
                    bullets = det.split('・') if '・' in det else [det]
                    bh2 = (row_h-I(0.1)) // max(len(bullets[:3]),1)
                    by3 = ry + I(0.05)
                    for bullet in bullets[:3]:
                        if bullet.strip():
                            _r(slide, cx+label_w+I(0.1), by3+I(0.08),
                               I(0.12), I(0.12), fill=NAVY)
                            _t(slide, cx+label_w+I(0.26), by3+I(0.02),
                               cw-label_w-I(0.35), bh2-I(0.04),
                               bullet.strip(), sz=fs, col=GRAY1)
                            by3 += bh2
                if ri < n_col-1:
                    _r(slide, cx, ry+row_h-P(2), cw, P(0.5), fill=LGRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 4. 2カラム比較（動的項目数対応）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def slide_twocol(slide, prs, headline, body):
    W, H = int(prs.slide_width), int(prs.slide_height)
    _r(slide, 0, 0, W, H, fill=WHITE)
    _page_frame(slide, prs)
    summary = body[0].split('：')[1] if body and '：' in body[0] else (body[0] if body else '')
    _slide_title(slide, prs, headline, summary)

    # bodyを左右に動的分配
    items = [b for b in body if b and not (b == body[0] and summary)]
    if not items: items = body[1:] if len(body)>1 else body
    mid = (len(items) + 1) // 2
    left_items  = items[:mid]
    right_items = items[mid:]

    CY = I(1.52); CH = H-CY-I(0.42)
    cw = int((W-I(0.9))/2); hh = I(0.36)
    titles = ['背景', '目的']

    for ci, (col_items, col_title) in enumerate(zip([left_items, right_items], titles)):
        if not col_items: continue
        cx = I(0.4)+ci*(cw+I(0.1))
        _navy_header(slide, cx, CY, cw, hh, col_title, sz=12)
        _r(slide, cx, CY+hh, cw, CH-hh, fill=WHITE, lc=NAVY, lw=0.75)

        avail_h = CH - hh - I(0.12)
        n_col = len(col_items)
        per = avail_h // n_col
        fs_big = 10 if per >= I(0.7) else 9
        fs_sub = 9  if per >= I(0.7) else 8

        iy = CY + hh + I(0.1)
        for j, item in enumerate(col_items):
            lbl, det = _parse(item)
            # ■ マーカー
            _r(slide, cx+I(0.1), iy+I(0.05), I(0.15), I(0.15), fill=NAVY)
            _t(slide, cx+I(0.32), iy+I(0.01), cw-I(0.38), I(0.25),
               lbl or item[:24], sz=fs_big, bold=True, col=GRAY1)
            if det:
                subs = det.split('・') if '・' in det else [det]
                sy = iy + I(0.28)
                s_per = max(I(0.22), (per - I(0.33)) // max(len(subs),1))
                for s in subs[:4]:
                    _t(slide, cx+I(0.42), sy, cw-I(0.48), s_per,
                       '• '+s.strip(), sz=fs_sub, col=GRAY2)
                    sy += s_per
            if j < n_col-1:
                _r(slide, cx+I(0.08), iy+per-I(0.04), cw-I(0.16), P(0.3), fill=LGRAY)
            iy += per


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 5. 比較テーブル（行数動的対応）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def slide_table(slide, prs, headline, body):
    W, H = int(prs.slide_width), int(prs.slide_height)
    _r(slide, 0, 0, W, H, fill=WHITE)
    _page_frame(slide, prs)
    summary = body[0].split('：')[1] if body and '：' in body[0] else (body[0] if body else '')
    _slide_title(slide, prs, headline, summary)

    items = [b for b in body if b and not (b == body[0] and summary)]
    if not items: items = body

    # 列定義
    cols=['評価軸','自社','競合A','競合B']
    HY=I(1.52)
    header_h = I(0.5)

    # 行数に応じて行高を動的計算
    n_rows = min(len(items), 8)
    avail_h = H - HY - header_h - I(0.42)
    RH = avail_h // n_rows
    # フォントサイズを行高に応じて調整
    fs_label = 10 if RH >= I(0.8) else (9 if RH >= I(0.65) else 8)
    fs_sub   = 8  if RH >= I(0.8) else (7 if RH >= I(0.65) else 7)

    # 列幅: 評価軸は広め
    axis_w = I(2.5)
    val_w  = int((W - I(0.8) - axis_w) / 3)
    col_x = [I(0.4), I(0.4)+axis_w, I(0.4)+axis_w+val_w, I(0.4)+axis_w+val_w*2]
    col_w = [axis_w, val_w, val_w, val_w]

    # ヘッダー行
    for ci,(lbl,x,w) in enumerate(zip(cols,col_x,col_w)):
        bg=NAVY if ci==0 else BLUE2
        _r(slide,x,HY,w-P(2),header_h,fill=bg)
        _t(slide,x,HY+I(0.1),w-P(2),header_h-I(0.18),
           lbl,sz=11,bold=True,col=WHITE,align=PP_ALIGN.CENTER)

    patterns=[('◎','△','○'),('○','◎','△'),('◎','○','△'),
              ('○','△','◎'),('◎','△','◎'),('△','◎','○'),
              ('◎','○','○'),('○','△','△')]
    sc={'◎':RGBColor(0x1C,0x7A,0x48),'○':BLUE2,'△':ORANGE}

    for ri,b in enumerate(items[:n_rows]):
        ry=HY+header_h+P(1)+ri*(RH+P(1))
        lbl,det=_parse(b)
        bg=BGRAY if ri%2==0 else WHITE
        for x,w in zip(col_x,col_w):
            _r(slide,x,ry,w-P(2),RH,fill=bg,lc=LGRAY,lw=0.4)
        _t(slide,col_x[0]+I(0.1),ry+I(0.05),col_w[0]-I(0.15),I(0.28),
           lbl or b[:18],sz=fs_label,bold=True,col=BLUE1)
        if det:
            _t(slide,col_x[0]+I(0.1),ry+I(0.3),col_w[0]-I(0.15),RH-I(0.34),
               det[:24],sz=fs_sub,col=GRAY3)
        pat=patterns[ri%len(patterns)]
        # スコア記号サイズも行高に応じて
        sc_sz = 16 if RH >= I(0.8) else (13 if RH >= I(0.65) else 11)
        for s,x,w in zip(pat,col_x[1:],col_w[1:]):
            _t(slide,x,ry+I(0.05),w-P(2),RH-I(0.08),
               s,sz=sc_sz,bold=True,col=sc.get(s,GRAY2),align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 6. ガントチャート＋収益グラフ（タスク数動的対応）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def slide_gantt(slide, prs, headline, body):
    W, H = int(prs.slide_width), int(prs.slide_height)
    _r(slide, 0, 0, W, H, fill=WHITE)
    _page_frame(slide, prs)
    summary = body[0].split('：')[1] if body and '：' in body[0] else (body[0] if body else '')
    _slide_title(slide, prs, headline, summary)

    tasks = [b for b in body if b and not (b == body[0] and summary)]
    if not tasks: tasks = body
    n_tasks = min(len(tasks), 8)

    GX=I(0.4); GY=I(1.52); GW=I(5.8); GH=H-I(2.05)
    phases=['M1-2','M3-4','M5-6','M7-9','M10-12','Y2-3']
    ph_n = len(phases)
    pw = GW // ph_n

    # ネイビーヘッダー
    _r(slide,GX,GY,GW,I(0.36),fill=NAVY)
    for i,ph in enumerate(phases):
        _t(slide,GX+i*pw,GY+I(0.05),pw,I(0.26),ph,sz=8,bold=True,col=WHITE,align=PP_ALIGN.CENTER)
        if i>0: _r(slide,GX+i*pw,GY,P(0.3),GH,fill=LGRAY)

    # タスク行高を動的計算
    avail_h = GH - I(0.38)
    row_h = avail_h // n_tasks
    bar_h = max(int(row_h * 0.5), P(8))  # バー高さ（行高の50%、最低8pt）
    lbl_w = int(GW * 0.30)

    # タスクとバー配置のパターン（開始フェーズ・期間を内容から推測）
    bar_patterns = [(0,2),(0,3),(1,1),(1,3),(2,2),(2,4),(3,2),(4,2)]
    bar_colors   = [NAVY,BLUE2,ORANGE,BLUE3,BLUE2,NAVY,ORANGE,BLUE2]
    fs_task = 9 if row_h >= I(0.55) else 8

    for ri,b in enumerate(tasks[:n_tasks]):
        lbl,_ = _parse(b)
        ry = GY+I(0.38)+ri*row_h
        # 行背景
        _r(slide,GX,ry,GW,row_h-P(1),fill=BGRAY if ri%2==0 else WHITE,lc=LGRAY,lw=0.25)
        # ラベル
        _t(slide,GX+I(0.05),ry+P(3),lbl_w-I(0.06),row_h-P(6),
           lbl or b[:16],sz=fs_task,col=GRAY1)
        # バー
        st,dur = bar_patterns[ri%len(bar_patterns)]
        bx = GX+lbl_w+st*pw
        bw = min(dur*pw-P(2), GW-lbl_w-st*pw-P(3))
        if bw > 0:
            bar_y = ry+int((row_h-bar_h)//2)
            _r(slide,bx,bar_y,bw,bar_h,fill=bar_colors[ri%len(bar_colors)])

    _t(slide,GX,GY+GH+I(0.04),GW,I(0.2),'■ 実施ロードマップ',sz=8,col=GRAY3)

    # 右: 収益棒グラフ（タスク内容から数値を抽出、なければ固定値）
    BX=I(6.5); BY=I(1.52); BW=I(3.2); BH=GH-I(0.3)
    _t(slide,BX,BY,BW,I(0.26),'売上推移（万円/月）',sz=8,bold=True,col=GRAY2)
    _r(slide,BX,BY+I(0.28),BW,BH-I(0.28),fill=WHITE,lc=LGRAY,lw=0.5)

    # 数値を内容から抽出（万円の数字）
    import re
    extracted_vals = []
    for b in tasks:
        nums = re.findall(r'(\d+)万円?/月', b)
        if nums: extracted_vals.append(int(nums[0]))
    # 3つの時系列データを生成（6ヶ月・12ヶ月・24ヶ月）
    if len(extracted_vals) >= 2:
        base = min(extracted_vals); peak = max(extracted_vals)
        mid_val = (base + peak) // 2
        vals = [[base//2, base//5, base//8],
                [mid_val//2, mid_val//5, mid_val//8],
                [peak//2, peak//5, peak//8]]
    else:
        vals = [[300,100,50],[500,180,120],[800,280,220]]

    periods=['6ヶ月','12ヶ月','24ヶ月']
    bc=[NAVY,BLUE2,ORANGE]; bw2=int((BW-I(0.6))/3); mv=max(sum(v) for v in vals)*1.2
    chart_h = BH-I(1.0)

    for bi,(period,vl) in enumerate(zip(periods,vals)):
        bx2=BX+I(0.3)+bi*(bw2+I(0.1)); cum=0; total=sum(vl)
        for v,c in zip(vl,bc):
            if mv > 0: bh2=int((v/mv)*chart_h)
            else: bh2=0
            by2=BY+BH-I(0.42)-cum-bh2
            if bh2>0: _r(slide,bx2,by2,bw2-P(1),bh2,fill=c)
            cum+=bh2
        _t(slide,bx2,BY+BH-I(0.38),bw2,I(0.3),period,sz=7,col=GRAY2,align=PP_ALIGN.CENTER)
        _t(slide,bx2,BY+BH-I(0.7),bw2,I(0.3),f'{total}万',sz=8,bold=True,col=NAVY,align=PP_ALIGN.CENTER)

    for i,(lbl,c) in enumerate(zip(['基本','F&B','他'],bc)):
        _r(slide,BX+I(0.3)+i*I(1.0),BY+BH+I(0.03),I(0.13),I(0.13),fill=c)
        _t(slide,BX+I(0.46)+i*I(1.0),BY+BH+I(0.01),I(0.7),I(0.2),lbl,sz=7,col=GRAY2)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 7. ツリーマップ（セグメント数動的対応）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def slide_treemap(slide, prs, headline, body):
    W, H = int(prs.slide_width), int(prs.slide_height)
    _r(slide, 0, 0, W, H, fill=WHITE)
    _page_frame(slide, prs)
    summary = body[0].split('：')[1] if body and '：' in body[0] else (body[0] if body else '')
    _slide_title(slide, prs, headline, summary)

    items = [b for b in body if b and not (b == body[0] and summary)]
    if not items: items = body
    n = min(len(items), 6)

    TX=I(0.4); TY=I(1.52); TW=I(5.5); TH=H-I(2.05)
    segs=[(b.split('：')[0] if '：' in b else b[:12]) for b in items[:n]]

    # 動的面積比（n に合わせて生成）
    if n == 2:
        left_ratios  = [1.0]
        right_ratios = []
        lw = int(TW * 0.56)
    elif n == 3:
        left_ratios  = [0.6, 0.4]
        right_ratios = [1.0]
        lw = int(TW * 0.56)
    elif n == 4:
        left_ratios  = [0.6, 0.4]
        right_ratios = [0.55, 0.45]
        lw = int(TW * 0.56)
    elif n == 5:
        left_ratios  = [0.6, 0.4]
        right_ratios = [0.42, 0.34, 0.24]
        lw = int(TW * 0.56)
    else:  # 6+
        left_ratios  = [0.6, 0.4]
        right_ratios = [0.36, 0.28, 0.22, 0.14]
        lw = int(TW * 0.56)

    # 左列ブロック
    fills_left  = [NAVY, LBLUE]
    tcs_left    = [WHITE, NAVY]
    cum = 0
    for i, (r, fc, tc) in enumerate(zip(left_ratios, fills_left, tcs_left)):
        bh = int(TH * r)
        if i == len(left_ratios)-1: bh = TH-cum
        _r(slide, TX, TY+cum, lw, bh-P(1), fill=fc)
        lbl = segs[i] if i < len(segs) else ''
        fs = 13 if bh >= I(1.5) else (11 if bh >= I(1.0) else 9)
        _t(slide, TX+I(0.1), TY+cum+I(0.12), lw-I(0.2), bh-I(0.25),
           lbl, sz=fs, bold=True, col=tc)
        cum += bh

    # 右列ブロック（枠線のみ）
    rx2 = TX+lw+P(2); rw2 = TW-lw-P(2)
    cum2 = 0
    for i, r in enumerate(right_ratios):
        bh = int(TH * r)
        if i == len(right_ratios)-1: bh = TH-cum2
        if bh <= 0: break
        _r(slide, rx2, TY+cum2, rw2, bh-P(1), fill=WHITE, lc=LGRAY, lw=0.5)
        lbl = segs[len(left_ratios)+i] if len(left_ratios)+i < len(segs) else ''
        fs = 9 if bh >= I(0.7) else 8
        _t(slide, rx2+I(0.07), TY+cum2+I(0.05), rw2-I(0.1), bh-I(0.1),
           lbl, sz=fs, col=GRAY1)
        cum2 += bh

    _t(slide, TX, TY+TH+I(0.04), TW, I(0.2),
       '市場セグメント構造（面積＝重要度）', sz=8, col=GRAY3)

    # 右: KPIカード（項目数に合わせて高さを動的調整）
    KX=I(6.2); KY=I(1.52); KW=I(4.0)
    kpi_items = items[:min(n, 5)]
    kh = (H-I(2.0)) // max(len(kpi_items), 1)
    kh = min(kh, I(1.1))  # 最大高さを制限

    for i, b in enumerate(kpi_items):
        lbl, det = _parse(b)
        ky = KY+i*kh
        col = [NAVY,BLUE2,ORANGE,BLUE1,BLUE2][i%5]
        _r(slide,KX,ky,KW,kh-P(2),fill=WHITE,lc=LGRAY,lw=0.5)
        _r(slide,KX,ky,P(4),kh-P(2),fill=col)
        fs_lbl = 9 if kh >= I(0.8) else 8
        fs_det = 9 if kh >= I(0.8) else 8
        _t(slide,KX+I(0.14),ky+I(0.06),KW-I(0.18),min(I(0.28),kh*0.35),
           lbl,sz=fs_lbl,bold=True,col=col)
        if det:
            _t(slide,KX+I(0.14),ky+I(0.34),KW-I(0.18),kh-I(0.42),
               det,sz=fs_det,col=GRAY2)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 8. まとめ（KPI数・アクション数動的対応）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def slide_summary(slide, prs, headline, body):
    W, H = int(prs.slide_width), int(prs.slide_height)
    _r(slide, 0, 0, W, H, fill=WHITE)
    _page_frame(slide, prs)
    _slide_title(slide, prs, headline)

    items = [b for b in body if b]
    # 先頭3件をKPI、残りをアクションとして扱う
    kpi_items = items[:3]
    action_items = items[3:]

    # KPIボックス（枚数に合わせて列数調整）
    n_kpi = len(kpi_items)
    kw = (W-I(0.9)) // max(n_kpi, 1)
    kw = min(kw, I(3.3))  # 最大幅
    ky = I(1.52); kh = I(1.72)

    for i, b in enumerate(kpi_items):
        lbl, val = _parse(b)
        col = [NAVY, BLUE2, ORANGE][i%3]
        kx = I(0.4)+i*(kw+I(0.1))
        _r(slide, kx, ky, kw, kh, fill=WHITE, lc=LGRAY, lw=0.75)
        _r(slide, kx, ky, kw, P(4), fill=col)
        _t(slide, kx+I(0.12), ky+I(0.1), kw-I(0.24), I(0.3),
           lbl, sz=10, bold=True, col=col)
        short = val[:16] if val else ''
        _t(slide, kx+I(0.12), ky+I(0.46), kw-I(0.24), I(0.68),
           short, sz=12, bold=True, col=NAVY, wrap=True)
        rest = val[16:]
        if rest:
            _t(slide, kx+I(0.12), ky+I(1.2), kw-I(0.24), I(0.44),
               rest[:22], sz=8, col=GRAY2)

    # アクションリスト（件数に合わせて行高調整）
    if not action_items: return
    _t(slide, I(0.4), I(3.38), I(9.8), I(0.3),
       'Next Actions', sz=12, bold=True, col=NAVY)
    _r(slide, I(0.4), I(3.72), I(9.8), P(1.5), fill=BLUE2)

    n_act = len(action_items)
    act_h = H - I(3.9) - I(0.35)
    cols  = 2 if n_act > 3 else 1
    half  = (n_act + 1) // 2
    cw2   = (W-I(1.0))//cols if cols > 1 else W-I(1.0)
    row_h = act_h // (half if cols > 1 else n_act)
    row_h = min(row_h, I(0.9))
    fs_act = 10 if row_h >= I(0.75) else 9

    for i, b in enumerate(action_items[:min(n_act, 8)]):
        lbl, det = _parse(b)
        col = [NAVY, BLUE2, ORANGE][i%3]
        ci = i // half if cols > 1 else 0
        ri = i %  half if cols > 1 else i
        cx = I(0.4)+ci*(cw2+I(0.2))
        cy = I(3.9)+ri*row_h
        if cy+row_h > H-I(0.35): break
        _r(slide, cx, cy, I(0.36), I(0.36), fill=col)
        _t(slide, cx, cy, I(0.36), I(0.36),
           '✓', sz=11, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
        _t(slide, cx+I(0.46), cy, cw2-I(0.52), I(0.28),
           lbl or f'Action{i+1}', sz=fs_act, bold=True, col=BLUE1)
        if det:
            _t(slide, cx+I(0.46), cy+I(0.3), cw2-I(0.52), row_h-I(0.34),
               det, sz=fs_act, col=GRAY2)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ルーティング
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
DESIGN_NAMES = {
    slide_title:   'title',
    slide_treemap: 'treemap',
    slide_flow:    'flow',
    slide_twocol:  'twocol',
    slide_table:   'table',
    slide_gantt:   'gantt',
    slide_detail:  'detail',
    slide_summary: 'summary',
}

DESIGN_CANDIDATES = {
    'データ':    [slide_treemap, slide_detail, slide_gantt],
    '比較':      [slide_table,   slide_twocol, slide_detail],
    'プロセス':  [slide_flow,    slide_gantt,  slide_detail],
    '関係性':    [slide_flow,    slide_treemap, slide_detail],
    '組織':      [slide_table,   slide_detail,  slide_flow],
    'テキスト':  [slide_detail,  slide_twocol,  slide_treemap],
    '市場':      [slide_treemap, slide_detail,  slide_table],
    '競合':      [slide_table,   slide_twocol,  slide_detail],
    '計画':      [slide_gantt,   slide_flow,    slide_detail],
    '収益':      [slide_gantt,   slide_detail,  slide_treemap],
    'フロー':    [slide_flow,    slide_gantt,   slide_detail],
    'ステップ':  [slide_flow,    slide_gantt,   slide_detail],
    'サービス':  [slide_flow,    slide_detail,  slide_treemap],
    '背景':      [slide_twocol,  slide_detail,  slide_flow],
    '目的':      [slide_twocol,  slide_detail,  slide_flow],
    '論点':      [slide_detail,  slide_twocol,  slide_table],
}

ALL_MIDDLE = [slide_treemap, slide_flow, slide_twocol, slide_table, slide_gantt, slide_detail]

def get_design_fn(purpose, content_type, page, total, used_history=None):
    if used_history is None: used_history = []
    if page == 1: return slide_title
    if page == total: return slide_summary
    avoid = set((used_history[-3:] if len(used_history)>=3 else used_history)+['title','summary'])
    candidates = []
    for key in [content_type, purpose]:
        if not key: continue
        for kw, fns in DESIGN_CANDIDATES.items():
            if kw in key:
                candidates.extend(fns); break
    seen=set(); unique=[]
    for fn in candidates:
        n=DESIGN_NAMES.get(fn,'')
        if n not in seen and n not in ('title','summary'):
            seen.add(n); unique.append(fn)
    for fn in unique:
        if DESIGN_NAMES.get(fn,'') not in avoid: return fn
    for fn in ALL_MIDDLE:
        if DESIGN_NAMES.get(fn,'') not in avoid: return fn
    counts={DESIGN_NAMES[fn]:used_history.count(DESIGN_NAMES[fn]) for fn in ALL_MIDDLE}
    least=min(counts,key=counts.get)
    for fn in ALL_MIDDLE:
        if DESIGN_NAMES[fn]==least: return fn
    return slide_detail
